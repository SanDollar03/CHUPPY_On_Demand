# app.py
# -*- coding: utf-8 -*-
import os
import re
import json
import time
import threading
import uuid
from collections import deque
from datetime import datetime
from typing import Dict, List, Tuple, Optional, Any

import requests
from flask import Flask, render_template, request, jsonify
from dotenv import load_dotenv

from docx import Document
from pypdf import PdfReader
from openpyxl import load_workbook
import xlrd
from pptx import Presentation


load_dotenv()

APP_TITLE = "CHUPPY RAG CONVERTER"
HEADER_MODEL_LABEL = "Model : ChatGPT 5.2"

API_BASE = (os.getenv("DIFY_API_BASE") or "").strip().rstrip("/")
API_KEY = (os.getenv("DIFY_API_KEY") or "").strip()
DATASET_API_BASE = (os.getenv("DIFY_DATASET_API_BASE") or API_BASE).strip().rstrip("/")
DATASET_API_KEY = (os.getenv("DIFY_DATASET_API_KEY") or API_KEY).strip()
DATASET_NAME_PREFIX = (os.getenv("DATASET_NAME_PREFIX") or "Chu_").strip()

EXPLORER_ROOT = os.path.normpath(
    os.getenv("CHUPPY_EXPLORER_ROOT") or r"\\172.27.23.54\disk1\Chuppy"
)
UPLOAD_MAX_FILES = 100
UPLOAD_MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024
EXPLORER_MAX_DEPTH = 5
UPLOAD_ALLOWED_DEPTH = 5

ALLOWED_EXTS = {
    ".txt", ".md", ".csv", ".json", ".log",
    ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
    ".py", ".js", ".css",
    ".docx", ".pdf",
    ".xlsx", ".xls", ".xlsm",
    ".ppt", ".pptx",
}
ALLOWED_EXTS_SORTED = tuple(sorted(ALLOWED_EXTS))
ALLOWED_EXTS_DISPLAY = ", ".join(ALLOWED_EXTS_SORTED)

MAX_INPUT_CHARS = 180_000
DEFAULT_CHUNK_SEP = "***"
REQ_TIMEOUT_SEC = 300
INDEXING_POLL_SEC = float(os.getenv("DIFY_INDEXING_POLL_SEC") or "2.0")
INDEXING_MAX_WAIT_SEC = int(os.getenv("DIFY_INDEXING_MAX_WAIT_SEC") or "900")
DIFY_MAX_SEG_TOKENS = int(os.getenv("DIFY_MAX_SEG_TOKENS") or "2000")

ONDEMAND_QUEUE_MAX_RETRIES = int(os.getenv("ONDEMAND_QUEUE_MAX_RETRIES") or "5")
ONDEMAND_QUEUE_HISTORY_LIMIT = int(os.getenv("ONDEMAND_QUEUE_HISTORY_LIMIT") or "300")
ONDEMAND_DATASET_CACHE_TTL_SEC = int(os.getenv("ONDEMAND_DATASET_CACHE_TTL_SEC") or "60")
ONDEMAND_DOCUMENT_CACHE_TTL_SEC = int(os.getenv("ONDEMAND_DOCUMENT_CACHE_TTL_SEC") or "60")
ONDEMAND_QUEUE_USER = (os.getenv("ONDEMAND_QUEUE_USER") or "rag_converter").strip() or "rag_converter"
ONDEMAND_QUEUE_STYLE = "rag_markdown"
ONDEMAND_QUEUE_CHUNK_SEP = DEFAULT_CHUNK_SEP
ONDEMAND_MONITOR_ENABLED = str(os.getenv("ONDEMAND_MONITOR_ENABLED") or "1").strip().lower() not in {"0", "false", "no", "off"}
ONDEMAND_MONITOR_INTERVAL_SEC = max(3.0, float(os.getenv("ONDEMAND_MONITOR_INTERVAL_SEC") or "15"))
ONDEMAND_SEEN_SIGNATURE_LIMIT = max(1000, int(os.getenv("ONDEMAND_SEEN_SIGNATURE_LIMIT") or "5000"))
DIFY_DELETE_HEALTH_CACHE_TTL_SEC = max(3, int(os.getenv("ONDEMAND_DIFY_DELETE_HEALTH_CACHE_TTL_SEC") or "10"))

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ONDEMAND_QUEUE_LOG_PATH = os.path.join(BASE_DIR, "ondemand_queue.log")
ONDEMAND_QUEUE_LOG_LOCK = threading.RLock()

def normalize_extension(ext: str) -> str:
    raw = str(ext or "").strip().lower()
    if not raw:
        return ""
    return raw if raw.startswith(".") else f".{raw}"


def is_allowed_extension(ext: str) -> bool:
    return normalize_extension(ext) in ALLOWED_EXTS


def is_allowed_filename(filename: str) -> bool:
    safe = sanitize_upload_filename(filename)
    if not safe:
        return False
    ext = os.path.splitext(safe)[1].lower()
    return is_allowed_extension(ext)


def get_allowed_extensions() -> List[str]:
    return list(ALLOWED_EXTS_SORTED)


def get_allowed_extensions_text() -> str:
    return ALLOWED_EXTS_DISPLAY


def get_upload_max_file_size_bytes() -> int:
    return int(UPLOAD_MAX_FILE_SIZE_BYTES)


def get_filestorage_size_bytes(file_storage: Any) -> int:
    try:
        size = int(getattr(file_storage, "content_length", 0) or 0)
        if size > 0:
            return size
    except Exception:
        pass

    stream = getattr(file_storage, "stream", None)
    if stream is None:
        return 0

    pos = None
    try:
        pos = stream.tell()
    except Exception:
        pos = None

    try:
        stream.seek(0, os.SEEK_END)
        end = int(stream.tell() or 0)
    except Exception:
        end = 0
    finally:
        if pos is not None:
            try:
                stream.seek(pos, os.SEEK_SET)
            except Exception:
                pass
        else:
            try:
                stream.seek(0, os.SEEK_SET)
            except Exception:
                pass
    return end



def ensure_queue_log_file() -> None:
    try:
        os.makedirs(os.path.dirname(ONDEMAND_QUEUE_LOG_PATH), exist_ok=True)
        need_header = (not os.path.exists(ONDEMAND_QUEUE_LOG_PATH)) or os.path.getsize(ONDEMAND_QUEUE_LOG_PATH) == 0
        if need_header:
            with open(ONDEMAND_QUEUE_LOG_PATH, "w", encoding="utf-8", newline="") as f:
                f.write("timestamp\tfolder_name\tfile_name\tstatus\n")
    except Exception:
        pass


def _compact_log_text(value: Any) -> str:
    return str(value or "").replace("\t", " ").replace("\r", " ").replace("\n", " ").strip()


def append_ondemand_queue_log(event: str, task: Optional[Dict[str, Any]] = None, **extra: Any) -> None:
    event_key = str(event or "").strip().lower()
    status = _compact_log_text(extra.get("status") or extra.get("log_status") or "")
    if not status:
        status = {
            "task_queued": "queued",
            "task_rejected": "rejected",
            "task_retry_reset": "retry_reset",
            "task_start": "running",
            "task_retry_wait": "retry_wait",
            "task_completed": "completed",
            "task_skipped": "skipped",
            "task_error": "error",
            "task_deleted": "deleted",
            "task_delete_restore": "queued",
            "markdown_cleanup": "markdown_deleted",
            "source_purged": "purged",
        }.get(event_key, "")
    if not status and isinstance(task, dict) and task:
        status = _compact_log_text(task.get("status") or task.get("result") or event_key or "queue")
    if not status:
        status = event_key or "queue"

    folder_rel_path = _compact_log_text(
        extra.get("folder_name")
        or extra.get("folder_rel_path")
        or ((task or {}).get("folder_rel_path") if isinstance(task, dict) else "")
    )
    file_name = _compact_log_text(
        extra.get("file_name")
        or extra.get("source_display_name")
        or ((task or {}).get("source_display_name") if isinstance(task, dict) else "")
        or ((task or {}).get("source_saved_name") if isinstance(task, dict) else "")
    )
    line = "\t".join([
        _compact_log_text(now_label()),
        folder_rel_path,
        file_name,
        status,
    ])

    try:
        with ONDEMAND_QUEUE_LOG_LOCK:
            with open(ONDEMAND_QUEUE_LOG_PATH, "a", encoding="utf-8", newline="") as f:
                f.write(line + "\n")
    except Exception:
        pass


def extract_text(path: str, knowledge_style: str = "rag_markdown") -> Tuple[str, Dict[str, str]]:
    ext = os.path.splitext(path)[1].lower()
    stat = os.stat(path)
    meta = {
        "filename": os.path.basename(path),
        "ext": ext,
        "size_bytes": str(stat.st_size),
        "mtime": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
    }

    if ext in {
        ".txt", ".md", ".csv", ".json", ".log",
        ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
        ".py", ".js", ".css",
    }:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), meta

    if ext == ".docx":
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                parts.append(t)
        return "\n".join(parts), meta

    if ext == ".pdf":
        return extract_pdf_like(path), meta

    if ext in {".xlsx", ".xlsm", ".xls"}:
        if knowledge_style == "rag_natural":
            text = extract_excel_as_markdown_tables(path, ext)
        else:
            text = extract_excel_as_row_records(path, ext)
        return text, meta

    if ext in {".ppt", ".pptx"}:
        return extract_ppt_like(path, ext), meta

    raise RuntimeError(f"未対応の拡張子です: {ext}")


def extract_pdf_like(path: str) -> str:
    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[PAGE {i+1}]\n{txt}")
    return "\n\n".join(parts)


def extract_excel_as_row_records(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_row_records(path)
    return extract_xlsx_like_as_row_records(path)


def extract_xlsx_like_as_row_records(path: str) -> str:
    wb = load_workbook(path, data_only=True, read_only=True)
    out: List[str] = []

    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("[EMPTY]")
            out.append("")
            continue

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_row_records(path: str) -> str:
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("[EMPTY]")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_excel_as_markdown_tables(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_markdown_tables(path)
    return extract_xlsx_like_as_markdown_tables(path)


def extract_xlsx_like_as_markdown_tables(path: str) -> str:
    max_rows_per_sheet = 200
    wb = load_workbook(path, data_only=True, read_only=True)

    out: List[str] = []
    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("(empty)")
            out.append("")
            continue

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= max_rows_per_sheet:
                out.append(f"(… {max_rows_per_sheet}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_markdown_tables(path: str) -> str:
    max_rows_per_sheet = 200
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("(empty)")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= max_rows_per_sheet:
                out.append(f"(… {max_rows_per_sheet}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_ppt_like(path: str, ext: str) -> str:
    try:
        prs = Presentation(path)
    except Exception:
        if ext == ".ppt":
            raise RuntimeError("`.ppt`（旧形式）は python-pptx で直接読めない場合があります。`.pptx` に変換して再実行してください。")
        raise RuntimeError("PowerPointの解析に失敗しました。ファイル破損または形式が想定外です。")

    parts: List[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text.append(t)

        txt = "\n".join(slide_text)
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[SLIDE {i+1}]\n{txt}")

    return "\n\n".join(parts)


def sanitize_header(cell) -> str:
    if cell is None:
        return ""
    s = str(cell).strip()
    s = re.sub(r"\s+", " ", s)
    return s


def normalize_pdf_like_text(s: str) -> str:
    lines = [ln.rstrip() for ln in s.splitlines()]
    out: List[str] = []
    buf = ""

    def flush():
        nonlocal buf
        if buf:
            out.append(buf)
            buf = ""

    for ln in lines:
        t = ln.strip("\u00a0 ").strip()
        if not t:
            flush()
            out.append("")
            continue
        if len(t) == 1:
            buf += t
        else:
            flush()
            out.append(t)
    flush()

    text = "\n".join(out)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def make_output_path(output_dir: str, rel_input_path: str) -> str:
    rel_dir = os.path.dirname(rel_input_path)
    base_name = os.path.splitext(os.path.basename(rel_input_path))[0]
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_name = f"{ts}_{base_name}.md"

    safe_dir = sanitize_relpath(rel_dir) if rel_dir else ""
    return os.path.join(output_dir, safe_dir, out_name)


def sanitize_relpath(p: str) -> str:
    if not p:
        return ""
    p = p.replace("..", "__")
    p = re.sub(r'[<>:"|?*]', "_", p)
    return p


def normalize_root_path(p: str) -> str:
    if not p:
        raise RuntimeError("ルートパスが未設定です。")
    return os.path.normcase(os.path.abspath(os.path.normpath(p)))


def make_rel_from_root(abs_path: str, root_dir: str) -> str:
    rel = os.path.relpath(abs_path, root_dir)
    if rel == ".":
        return ""
    return rel.replace("\\", "/")


def path_depth_from_rel(rel_path: str) -> int:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    if not rel:
        return 0
    return len([p for p in rel.split("/") if p])


def path_depth_from_root(abs_path: str, root_dir: str) -> int:
    rel = make_rel_from_root(abs_path, root_dir)
    return path_depth_from_rel(rel)


def resolve_explorer_path(root_dir: str, rel_path: str) -> str:
    root_norm = normalize_root_path(root_dir)

    rel = (rel_path or "").strip().replace("/", os.sep).replace("\\", os.sep)
    rel = rel.lstrip(os.sep)

    candidate = os.path.normpath(os.path.join(root_dir, rel))
    cand_norm = normalize_root_path(candidate)

    if cand_norm != root_norm and not cand_norm.startswith(root_norm + os.sep):
        raise RuntimeError("許可されていないパスです。")

    return candidate


def sanitize_upload_filename(name: str) -> str:
    original = os.path.basename((name or "").replace("\x00", "").strip())
    if not original or original in {".", ".."}:
        return ""

    safe = re.sub(r"[\x00-\x1f]", "", original)
    safe = safe.replace("/", "_").replace("\\", "_")
    safe = re.sub(r'[:*?"<>|]', "_", safe)
    safe = safe.rstrip(" .")

    if not safe or safe in {".", ".."}:
        return ""
    return safe


def add_upload_timestamp_prefix(filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        return ""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{ts}_{safe}"


def build_unique_upload_path(target_dir: str, filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        raise RuntimeError("使用できないファイル名です。")

    base, ext = os.path.splitext(safe)
    candidate = os.path.join(target_dir, safe)
    seq = 1
    while os.path.exists(candidate):
        candidate = os.path.join(target_dir, f"{base}_{seq:02d}{ext}")
        seq += 1
        if seq > 9999:
            raise RuntimeError("同名ファイルが多すぎるため保存できません。")
    return candidate


def matches_explorer_level_rule(depth: int, name: str) -> bool:
    nm = str(name or "")
    if depth == 1:
        return len(nm) == 1
    if depth == 2:
        return len(nm) == 2
    if depth == 4:
        return nm == "元データ"
    return True


def list_visible_child_dir_names(abs_dir: str, root_dir: str) -> List[str]:
    parent_depth = path_depth_from_root(abs_dir, root_dir)
    next_depth = parent_depth + 1
    if next_depth > EXPLORER_MAX_DEPTH:
        return []

    names: List[str] = []
    try:
        for name in os.listdir(abs_dir):
            full = os.path.join(abs_dir, name)
            if not os.path.isdir(full):
                continue
            if not matches_explorer_level_rule(next_depth, name):
                continue
            names.append(name)
    except Exception:
        return []

    names.sort(key=lambda x: x.lower())
    return names


def dir_has_child_dirs(abs_dir: str) -> bool:
    try:
        for name in os.listdir(abs_dir):
            full = os.path.join(abs_dir, name)
            if os.path.isdir(full):
                return True
    except Exception:
        return False
    return False


def dir_has_visible_child_dirs(abs_dir: str, root_dir: str) -> bool:
    return bool(list_visible_child_dir_names(abs_dir, root_dir))


def compute_visible_tree_stats(abs_dir: str, root_dir: str, cache: Optional[Dict[str, Dict[str, int]]] = None) -> Dict[str, int]:
    cache = cache if cache is not None else {}
    key = normalize_root_path(abs_dir)
    if key in cache:
        return cache[key]

    depth = path_depth_from_root(abs_dir, root_dir)
    child_names = list_visible_child_dir_names(abs_dir, root_dir)

    file_count = 0
    total_size_bytes = 0

    if depth >= EXPLORER_MAX_DEPTH or not child_names:
        try:
            for name in os.listdir(abs_dir):
                full = os.path.join(abs_dir, name)
                if not os.path.isfile(full):
                    continue
                ext = os.path.splitext(name)[1].lower()
                if not is_allowed_extension(ext):
                    continue
                try:
                    st = os.stat(full)
                except Exception:
                    continue
                file_count += 1
                total_size_bytes += int(st.st_size or 0)
        except Exception:
            pass
    else:
        for name in child_names:
            full = os.path.join(abs_dir, name)
            child_stats = compute_visible_tree_stats(full, root_dir, cache)
            file_count += int(child_stats.get("file_count") or 0)
            total_size_bytes += int(child_stats.get("total_size_bytes") or 0)

    out = {
        "file_count": file_count,
        "total_size_bytes": total_size_bytes,
    }
    cache[key] = out
    return out


def build_dir_info(abs_dir: str, root_dir: str, stats_cache: Optional[Dict[str, Dict[str, int]]] = None) -> Dict[str, Any]:
    depth = path_depth_from_root(abs_dir, root_dir)
    stats = compute_visible_tree_stats(abs_dir, root_dir, stats_cache)
    return {
        "name": os.path.basename(abs_dir.rstrip("\\/")) or abs_dir,
        "path": make_rel_from_root(abs_dir, root_dir),
        "abs_path": abs_dir,
        "depth": depth,
        "can_upload": depth == UPLOAD_ALLOWED_DEPTH,
        "has_children": dir_has_visible_child_dirs(abs_dir, root_dir),
        "file_count": int(stats.get("file_count") or 0),
        "total_size_bytes": int(stats.get("total_size_bytes") or 0),
    }


def list_explorer_dir(abs_dir: str, root_dir: str) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    dirs: List[Dict[str, Any]] = []
    files: List[Dict[str, Any]] = []
    stats_cache: Dict[str, Dict[str, int]] = {}

    current_depth = path_depth_from_root(abs_dir, root_dir)

    for name in sorted(os.listdir(abs_dir), key=lambda x: x.lower()):
        full = os.path.join(abs_dir, name)
        try:
            st = os.stat(full)
        except Exception:
            continue

        item = {
            "name": name,
            "path": make_rel_from_root(full, root_dir),
            "mtime": datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
        }

        if os.path.isdir(full):
            depth = path_depth_from_root(full, root_dir)
            if not matches_explorer_level_rule(depth, name):
                continue

            dir_info = build_dir_info(full, root_dir, stats_cache)
            item["type"] = "dir"
            item["depth"] = depth
            item["can_upload"] = depth == UPLOAD_ALLOWED_DEPTH
            item["has_children"] = False if depth >= EXPLORER_MAX_DEPTH else bool(dir_info.get("has_children"))
            item["expandable"] = depth < EXPLORER_MAX_DEPTH and bool(dir_info.get("has_children"))
            item["file_count"] = int(dir_info.get("file_count") or 0)
            item["total_size_bytes"] = int(dir_info.get("total_size_bytes") or 0)
            dirs.append(item)
        else:
            ext = os.path.splitext(name)[1].lower()
            if not is_allowed_extension(ext):
                continue
            item["type"] = "file"
            item["size_bytes"] = st.st_size
            item["depth"] = current_depth
            files.append(item)

    return dirs, files


def get_cached_dataset_document_name_keys_any(dataset_id: str) -> set:
    ds_id = str(dataset_id or "").strip()
    if not ds_id:
        return set()
    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(ds_id) or {}
        return set(entry.get("keys") or set())


def get_queue_registration_status_class(task_status: str) -> str:
    st = str(task_status or "").strip().lower()
    if st in {"queued", "running"}:
        return "registering"
    if st == "error":
        return "error"
    if st in {"completed", "skipped"}:
        return "registered"
    return ""


def build_registration_status_payload(
    code: str,
    label: str,
    detail: str = "",
    queue_status: str = "",
    queue_stage: str = "",
    queue_order: Optional[int] = None,
) -> Dict[str, Any]:
    return {
        "registration_status": code,
        "registration_status_label": label,
        "registration_status_detail": detail or "",
        "registration_queue_status": queue_status or "",
        "registration_queue_stage": queue_stage or "",
        "registration_queue_order": int(queue_order) if isinstance(queue_order, int) else (int(queue_order) if str(queue_order).isdigit() else None),
    }


def resolve_explorer_file_registration_status(
    folder_rel_path: str,
    file_item: Dict[str, Any],
    latest_task: Optional[Dict[str, Any]],
) -> Dict[str, Any]:
    file_name = str((file_item or {}).get("name") or "")
    queue_status = str((latest_task or {}).get("status") or "")
    queue_stage = str((latest_task or {}).get("stage") or "")
    queue_order = (latest_task or {}).get("queue_order")
    latest_error = str((latest_task or {}).get("last_error") or "")
    latest_message = str((latest_task or {}).get("message") or "")

    if queue_status in {"queued", "running"}:
        if queue_status == "running":
            detail = latest_message or "現在Markdown変換処理を実行しています。"
        else:
            detail = latest_message or "アップロード待ちキューで順番待ちです。"
        return build_registration_status_payload(
            code="registering",
            label="変換中",
            detail=detail,
            queue_status=queue_status,
            queue_stage=queue_stage,
            queue_order=queue_order,
        )

    markdown_abs_path = ""
    try:
        base_md_abs_path, _, _ = build_ondemand_markdown_path(folder_rel_path, file_name)
        markdown_abs_path = base_md_abs_path
    except Exception:
        markdown_abs_path = ""

    md_exists = bool(markdown_abs_path) and os.path.exists(markdown_abs_path)
    if md_exists:
        detail = "Markdownへ変換済みです。"
        if queue_status in {"completed", "skipped"} and latest_message:
            detail = latest_message
        return build_registration_status_payload(
            code="registered",
            label="変換済",
            detail=detail,
            queue_status=queue_status,
            queue_stage=queue_stage,
            queue_order=queue_order,
        )

    if queue_status == "error":
        return build_registration_status_payload(
            code="error",
            label="エラー",
            detail=latest_error or latest_message or "Markdown変換時にエラーが発生しました。",
            queue_status=queue_status,
            queue_stage=queue_stage,
            queue_order=queue_order,
        )

    return build_registration_status_payload(
        code="unregistered",
        label="未変換",
        detail="Markdown未変換です。",
        queue_status=queue_status,
        queue_stage=queue_stage,
        queue_order=queue_order,
    )


def enrich_explorer_files_with_registration_status(
    folder_rel_path: str,
    files: List[Dict[str, Any]],
    queue_manager: Optional[Any] = None,
    force: bool = False,
) -> List[Dict[str, Any]]:
    if not files:
        return []

    folder_rel = str(folder_rel_path or "").strip().replace("\\", "/").strip("/")

    source_rel_paths = [str((it or {}).get("path") or "") for it in (files or []) if (it or {}).get("path")]
    latest_task_map: Dict[str, Dict[str, Any]] = {}
    if queue_manager is not None:
        latest_task_map = queue_manager.get_latest_task_snapshots_by_source_rel_paths(source_rel_paths)

    out: List[Dict[str, Any]] = []
    for file_item in (files or []):
        item = dict(file_item or {})
        rel_key = normalize_name_key(str(item.get("path") or "").replace("\\", "/"))
        latest_task = latest_task_map.get(rel_key)
        status_payload = resolve_explorer_file_registration_status(
            folder_rel_path=folder_rel,
            file_item=item,
            latest_task=latest_task,
        )
        item.update(status_payload)
        out.append(item)
    return out


def sse_event(event: str, data: Dict[str, Any]) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


def format_size_bytes(size_bytes: Any) -> str:
    try:
        value = float(size_bytes or 0)
    except Exception:
        value = 0.0
    units = ["B", "KB", "MB", "GB", "TB"]
    idx = 0
    while value >= 1024 and idx < len(units) - 1:
        value /= 1024.0
        idx += 1
    if idx == 0 or value >= 10:
        return f"{value:.0f} {units[idx]}"
    return f"{value:.1f} {units[idx]}"


def safe_err(msg: str) -> str:
    if not msg:
        return "不明なエラー"
    msg = re.sub(r"(app-[A-Za-z0-9_\-]{10,})", "app-***REDACTED***", msg)
    msg = re.sub(r"(Bearer\s+)[A-Za-z0-9_\-\.=]+", r"\1***REDACTED***", msg, flags=re.IGNORECASE)
    msg = re.sub(r"https?://[^\s]+", "[URL_REDACTED]", msg)
    return msg[:700]


# -----------------------
# Chat conversion (/chat-messages)
# -----------------------

def convert_via_dify_chat_messages_secure(
    api_base: str,
    api_key: str,
    user: str,
    source_path: str,
    source_meta: Dict[str, str],
    text: str,
    knowledge_style: str,
    chunk_sep: str,
) -> str:
    url = f"{api_base}/chat-messages"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    instruction = build_rag_instruction(
        source_path=source_path,
        source_meta=source_meta,
        knowledge_style=knowledge_style,
        chunk_sep=chunk_sep,
    )

    query = (
        instruction
        + "\n\n===== SOURCE TEXT BEGIN =====\n"
        + text
        + "\n===== SOURCE TEXT END =====\n"
    )

    payload = {
        "inputs": {},
        "query": query,
        "response_mode": "blocking",
        "user": user,
    }

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=REQ_TIMEOUT_SEC)
    except requests.RequestException:
        raise RuntimeError("API通信に失敗しました（ネットワーク/タイムアウト）。")

    if r.status_code >= 400:
        raise RuntimeError(f"APIエラー（HTTP {r.status_code}）: {safe_err(r.text)}")

    try:
        data = r.json()
    except Exception:
        raise RuntimeError("APIレスポンスの解析に失敗しました。")

    answer = data.get("answer")
    if not answer or not isinstance(answer, str):
        raise RuntimeError("APIレスポンスが想定外です（answerがありません）。")

    return answer.strip() + "\n"


def build_rag_instruction(source_path: str, source_meta: Dict[str, str], knowledge_style: str, chunk_sep: str) -> str:
    meta_lines = "\n".join([f"- {k}: {v}" for k, v in source_meta.items()])
    ext = (source_meta.get("ext") or "").lower()

    first_chunk_rule = f"""
        # 最初のチャンク（必須）
        - 出力の最初のチャンクは必ず「全体構成（目次/分類）」にする。
        - 形式例：
        - 見出し: 「## 全体構成（目次/分類）」
        - 次の1文: 「このチャンクでは文書全体の構成（目次）と分類方針を示す。」
        - 続けて、章立て（大カテゴリ）と、その中で扱う内容の要約を箇条書きで書く。
        - そのチャンクの末尾に必ず「{chunk_sep}」を単独行で置く。
        """

    excel_rules = ""
    if ext in {".xlsx", ".xls", ".xlsm"} and knowledge_style != "rag_natural":
        excel_rules = f"""
        # Excel特別ルール（標準/FAQ用）
        - 入力には [HEADER] と [ROW] が含まれる。
        - 出力は「データ行（[ROW]）1つにつき、必ずチャンク1つ」にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - [ROW]を統合しない。行同士をまとめない。
        """

    if knowledge_style == "rag_natural":
        style_block = f"""
        出力はMarkdownで「RAG向けMarkdown（自然言語）」として整形する。

        # 手順（必須）
        1) まず文書全体の構成を把握し、上位の章立て（大カテゴリ）を作る。
        2) 次に、人間が指示を出すような自然文で各チャンクの目的を宣言してから本文を置く。
        3) チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    elif knowledge_style == "faq":
        style_block = f"""
        出力はMarkdownで、FAQ形式にする。
        - 質問は具体的に、回答は短く「結論→根拠→例」の順にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    else:
        style_block = f"""
        出力はMarkdownで、RAGに最適化したナレッジへ整形する。
        - 文は「主語 + 述語」でできるだけ明確にする。
        - 検索されやすいキーワード（固有名詞/手順名/条件/例外/閾値）を含める。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - 情報を省略しない（重複は統合可）。
        """

    return f"""
        あなたは「社内RAG用ナレッジ整形AI」である。
        入力された文章を、検索精度が最大化するMarkdownへ再構成する。

        # 変換対象ファイル
        - path: {source_path}
        - meta:
        {meta_lines}

        # 絶対ルール
        - 出力は「変換後Markdown本文のみ」とする（前置き/解説/謝罪/注釈は禁止）。
        - 原文が曖昧な場合は「〜である可能性がある」等で補い、捏造しない。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。

        {first_chunk_rule}

        {excel_rules}

        # スタイル
        {style_block}
        """.strip()


# -----------------------
# Markdown metadata + chunk analysis
# -----------------------

def _yaml_quote(v: str) -> str:
    s = "" if v is None else str(v)
    s = s.replace("\\", "\\\\").replace('"', '\\"')
    return f'"{s}"'


def attach_source_metadata(md: str, source_relpath: str, source_abspath: str, source_meta: Dict[str, str]) -> str:
    fm = {
        "source_relpath": source_relpath,
        "source_abspath": os.path.abspath(source_abspath),
        "source_filename": source_meta.get("filename") or os.path.basename(source_abspath),
        "source_ext": source_meta.get("ext") or os.path.splitext(source_abspath)[1].lower(),
        "source_size_bytes": source_meta.get("size_bytes") or "",
        "source_mtime": source_meta.get("mtime") or "",
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }

    lines = ["---"]
    for k, v in fm.items():
        lines.append(f"{k}: {_yaml_quote(v)}")
    lines.append("---")
    lines.append("")

    body = (md or "").lstrip("\ufeff\n\r ")
    return "\n".join(lines) + body


def normalize_chunk_sep_lines(md: str, chunk_sep: str) -> str:
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP
    lines = []
    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            lines.append(sep)
        else:
            lines.append(ln.rstrip("\r"))
    out = "\n".join(lines).strip()
    return out + "\n"


def split_chunks(md: str, chunk_sep: str) -> List[str]:
    chunks: List[str] = []
    buf: List[str] = []
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            txt = "\n".join(buf).strip()
            if txt:
                chunks.append(txt)
            buf = []
        else:
            buf.append(ln)

    last = "\n".join(buf).strip()
    if last:
        chunks.append(last)
    return chunks


def estimate_tokens(text: str) -> int:
    if not text:
        return 0
    total = len(text)
    if total <= 0:
        return 0

    ascii_cnt = sum(1 for ch in text if ord(ch) < 128)
    ascii_ratio = ascii_cnt / total

    chars_per_token = 3.0 if ascii_ratio >= 0.60 else 1.6
    est = int(total / chars_per_token) + 1
    return max(1, est)


def analyze_chunks_for_dify(markdown: str, chunk_sep: str) -> Dict[str, Any]:
    chunks = split_chunks(markdown, chunk_sep)
    lens = [estimate_tokens(c) for c in chunks] if chunks else []

    if not lens:
        return {
            "chunks": 0,
            "chunk_tokens_max": 0,
            "dify_max_tokens": min(1000, DIFY_MAX_SEG_TOKENS),
        }

    max_tok = max(lens)

    target = max_tok + 32
    target = max(200, target)
    target = min(DIFY_MAX_SEG_TOKENS, target)

    return {
        "chunks": len(chunks),
        "chunk_tokens_max": max_tok,
        "dify_max_tokens": target,
    }


# -----------------------
# Dify Knowledge API
# -----------------------

def dify_headers(api_key: str) -> Dict[str, str]:
    return {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }


def dify_list_datasets(api_base: str, api_key: str, prefix: str, limit: int = 100) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    page = 1

    while True:
        url = f"{api_base}/datasets?page={page}&limit={limit}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"datasets取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []

        for it in items:
            name = (it.get("name") or "").strip()
            if prefix and not name.startswith(prefix):
                continue
            did = (it.get("id") or "").strip()
            if did and name:
                out.append({"id": did, "name": name})

        has_more = bool(data.get("has_more"))
        if not has_more:
            break

        page += 1
        if page > 200:
            break

    return out


def dify_get_dataset_detail(api_base: str, api_key: str, dataset_id: str) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code == 405:
        return {
            "id": dataset_id,
            "name": "",
            "_note": "datasets/{id} が GET 非対応のため、詳細は省略しました。",
        }
    if r.status_code >= 400:
        raise RuntimeError(f"dataset詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_list_documents_all(
    api_base: str,
    api_key: str,
    dataset_id: str,
    keyword: str = "",
    limit: int = 100,
) -> Tuple[List[Dict[str, Any]], int]:
    items_out: List[Dict[str, Any]] = []
    page = 1
    total = 0

    while True:
        qs = f"page={page}&limit={limit}"
        if keyword:
            qs += "&keyword=" + requests.utils.quote(keyword)

        url = f"{api_base}/datasets/{dataset_id}/documents?{qs}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"documents取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []
        total = int(data.get("total") or total or 0)

        for it in items:
            if isinstance(it, dict):
                items_out.append(it)

        if not bool(data.get("has_more")):
            break

        page += 1
        if page > 200:
            break

    return items_out, total


def dify_get_document_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    metadata: str = "without",
) -> Dict[str, Any]:
    meta = metadata.strip() if metadata else "without"
    if meta not in {"all", "only", "without"}:
        meta = "without"

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}?metadata={meta}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"document詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_delete_document(dataset_id: str, document_id: str) -> None:
    ds_id = str(dataset_id or "").strip()
    doc_id = str(document_id or "").strip()
    if not ds_id or not doc_id:
        raise RuntimeError("削除対象documentを判定できません。")

    url = f"{DATASET_API_BASE}/datasets/{ds_id}/documents/{doc_id}"
    r = requests.delete(url, headers={"Authorization": f"Bearer {DATASET_API_KEY}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code not in {200, 202, 204}:
        raise RuntimeError(f"document削除失敗（HTTP {r.status_code}）: {safe_err(r.text)}")


def dify_list_segments_page(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    page: int = 1,
    limit: int = 20,
    keyword: str = "",
    status: str = "",
) -> Dict[str, Any]:
    qs = f"page={page}&limit={limit}"
    if keyword:
        qs += "&keyword=" + requests.utils.quote(keyword)
    if status:
        qs += "&status=" + requests.utils.quote(status)

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments?{qs}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segments取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return {
        "items": data.get("data") or [],
        "has_more": bool(data.get("has_more")),
        "total": int(data.get("total") or 0),
        "page": int(data.get("page") or page),
        "limit": int(data.get("limit") or limit),
    }


def dify_get_segment_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    segment_id: str,
) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments/{segment_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segment詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_create_document_by_text(
    dataset_id: str,
    name: str,
    text: str,
    chunk_sep: str,
    dify_max_tokens: int,
    search_method: str = "hybrid_search",
) -> Tuple[str, str]:
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/document/create-by-text"

    payload: Dict[str, Any] = {
        "name": name,
        "text": text,
        "indexing_technique": "high_quality",
        "doc_form": "text_model",
        "process_rule": {
            "mode": "custom",
            "rules": {
                "pre_processing_rules": [
                    {"id": "remove_extra_spaces", "enabled": True},
                    {"id": "remove_urls_emails", "enabled": True},
                ],
                "segmentation": {
                    "separator": chunk_sep,
                    "max_tokens": int(dify_max_tokens),
                },
            },
        },
        "retrieval_model": {
            "search_method": search_method,
            "reranking_enable": False,
            "top_k": 5,
            "score_threshold_enabled": False,
        },
    }

    r = requests.post(url, headers=dify_headers(DATASET_API_KEY), json=payload, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"create-by-text 失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    doc = data.get("document") or {}
    doc_id = (doc.get("id") or "").strip()
    batch = (data.get("batch") or "").strip()

    if not doc_id or not batch:
        raise RuntimeError("create-by-text レスポンスが想定外です（document.id / batch がありません）。")

    return doc_id, batch


def dify_get_indexing_status(dataset_id: str, batch: str) -> List[Dict[str, Any]]:
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/documents/{batch}/indexing-status"
    r = requests.get(url, headers={"Authorization": f"Bearer {DATASET_API_KEY}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"indexing-status 取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return data.get("data") or []


def register_markdown_to_dify(dataset_id: str, doc_name: str, markdown: str, chunk_sep: str) -> Dict[str, Any]:
    stats = analyze_chunks_for_dify(markdown, chunk_sep)
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    doc_id, batch = dify_create_document_by_text(
        dataset_id=dataset_id,
        name=doc_name,
        text=markdown,
        chunk_sep=sep,
        dify_max_tokens=int(stats["dify_max_tokens"]),
        search_method="hybrid_search",
    )

    return {
        "doc_id": doc_id,
        "batch": batch,
        "chunk_sep": sep,
        "chunks": stats["chunks"],
        "chunk_tokens_max": stats["chunk_tokens_max"],
        "dify_max_tokens": stats["dify_max_tokens"],
        "search_method": "hybrid_search",
    }


def iter_indexing_status(dataset_id: str, batch: str, doc_id: str):
    start = time.time()
    last_key = None

    while True:
        if time.time() - start > INDEXING_MAX_WAIT_SEC:
            raise RuntimeError("ナレッジ埋め込みがタイムアウトしました。")

        items = dify_get_indexing_status(dataset_id, batch)

        target = None
        for it in items:
            if (it.get("id") or "").strip() == doc_id:
                target = it
                break

        if not target:
            time.sleep(INDEXING_POLL_SEC)
            continue

        st = (target.get("indexing_status") or "").strip()
        completed = int(target.get("completed_segments") or 0)
        total = int(target.get("total_segments") or 0)
        err = target.get("error")

        key = f"{st}:{completed}/{total}:{err}"
        if key != last_key:
            last_key = key
            terminal = st.lower() in {"completed", "error", "failed", "stopped"}
            yield {
                "indexing_status": st,
                "completed_segments": completed,
                "total_segments": total,
                "error": err,
                "terminal": terminal,
            }

            if terminal:
                return

        time.sleep(INDEXING_POLL_SEC)



_DATASET_CACHE_LOCK = threading.RLock()
_DATASET_CACHE: Dict[str, Any] = {
    "ts": 0.0,
    "items": [],
}

_DOCUMENT_CACHE_LOCK = threading.RLock()
_DOCUMENT_NAME_CACHE: Dict[str, Dict[str, Any]] = {}

_DIFY_DELETE_HEALTH_LOCK = threading.RLock()
_DIFY_DELETE_HEALTH_CACHE: Dict[str, Any] = {
    "ts": 0.0,
    "ok": False,
    "reason": "",
}


def now_label() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def normalize_name_key(name: str) -> str:
    return str(name or "").strip().lower()


def get_datasets_cached(force: bool = False) -> List[Dict[str, Any]]:
    with _DATASET_CACHE_LOCK:
        cache_ts = float(_DATASET_CACHE.get("ts") or 0.0)
        cache_items = list(_DATASET_CACHE.get("items") or [])
        fresh = (time.time() - cache_ts) <= ONDEMAND_DATASET_CACHE_TTL_SEC
        if cache_items and fresh and not force:
            return cache_items

    items = dify_list_datasets(
        api_base=DATASET_API_BASE,
        api_key=DATASET_API_KEY,
        prefix=DATASET_NAME_PREFIX,
        limit=100,
    )

    with _DATASET_CACHE_LOCK:
        _DATASET_CACHE["ts"] = time.time()
        _DATASET_CACHE["items"] = list(items or [])
        return list(_DATASET_CACHE["items"])


def find_dataset_by_name(dataset_name: str) -> Optional[Dict[str, Any]]:
    if not dataset_name:
        return None

    key = normalize_name_key(dataset_name)
    items = get_datasets_cached(force=False)
    for it in items:
        if normalize_name_key(it.get("name")) == key:
            return dict(it)
    return None


def get_dataset_document_name_keys_cached(dataset_id: str, force: bool = False) -> set:
    dataset_id = (dataset_id or "").strip()
    if not dataset_id:
        return set()

    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(dataset_id) or {}
        cache_ts = float(entry.get("ts") or 0.0)
        cache_keys = set(entry.get("keys") or set())
        fresh = (time.time() - cache_ts) <= ONDEMAND_DOCUMENT_CACHE_TTL_SEC
        if cache_keys and fresh and not force:
            return cache_keys

    items, _ = dify_list_documents_all(
        api_base=DATASET_API_BASE,
        api_key=DATASET_API_KEY,
        dataset_id=dataset_id,
        keyword="",
        limit=100,
    )
    keys = set()
    for it in (items or []):
        name = (it or {}).get("name")
        if not name:
            continue
        status = normalize_name_key((it or {}).get("indexing_status") or (it or {}).get("display_status") or "")
        if status in {"error", "failed", "stopped"}:
            continue
        keys.add(normalize_name_key(name))

    with _DOCUMENT_CACHE_LOCK:
        _DOCUMENT_NAME_CACHE[dataset_id] = {
            "ts": time.time(),
            "keys": set(keys),
        }
        return set(keys)


def find_dataset_documents_by_name(dataset_id: str, doc_name: str) -> List[Dict[str, Any]]:
    ds_id = str(dataset_id or "").strip()
    nm = str(doc_name or "").strip()
    if not ds_id or not nm:
        return []

    items, _ = dify_list_documents_all(
        api_base=DATASET_API_BASE,
        api_key=DATASET_API_KEY,
        dataset_id=ds_id,
        keyword=nm,
        limit=100,
    )
    name_key = normalize_name_key(nm)
    out: List[Dict[str, Any]] = []
    for it in (items or []):
        if normalize_name_key((it or {}).get("name")) != name_key:
            continue
        doc_id = str((it or {}).get("id") or "").strip()
        if not doc_id:
            continue
        out.append({
            "id": doc_id,
            "name": (it or {}).get("name") or "",
            "indexing_status": (it or {}).get("indexing_status") or (it or {}).get("display_status") or "",
        })
    return out


def invalidate_dataset_document_cache(dataset_id: str) -> None:
    ds_id = str(dataset_id or "").strip()
    if not ds_id:
        return
    with _DOCUMENT_CACHE_LOCK:
        _DOCUMENT_NAME_CACHE.pop(ds_id, None)


def get_dify_delete_availability(force: bool = False) -> Tuple[bool, str]:
    return True, ""


def dataset_document_exists_by_name(dataset_id: str, doc_name: str, force: bool = False) -> bool:
    return normalize_name_key(doc_name) in get_dataset_document_name_keys_cached(dataset_id, force=force)


def remember_dataset_document_name(dataset_id: str, doc_name: str) -> None:
    dataset_id = (dataset_id or "").strip()
    if not dataset_id or not doc_name:
        return

    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(dataset_id) or {"ts": time.time(), "keys": set()}
        keys = set(entry.get("keys") or set())
        keys.add(normalize_name_key(doc_name))
        entry["keys"] = keys
        entry["ts"] = time.time()
        _DOCUMENT_NAME_CACHE[dataset_id] = entry


def forget_dataset_document_name(dataset_id: str, doc_name: str) -> None:
    dataset_id = (dataset_id or "").strip()
    if not dataset_id or not doc_name:
        return

    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(dataset_id) or {"ts": time.time(), "keys": set()}
        keys = set(entry.get("keys") or set())
        keys.discard(normalize_name_key(doc_name))
        entry["keys"] = keys
        entry["ts"] = time.time()
        _DOCUMENT_NAME_CACHE[dataset_id] = entry


def build_ondemand_dataset_name(rel_path: str) -> str:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    if len(parts) != UPLOAD_ALLOWED_DEPTH:
        return ""

    filtered = [p for p in parts if p != "元データ"]
    if not filtered:
        return ""

    return DATASET_NAME_PREFIX + "_".join(filtered)


def build_ondemand_markdown_path(folder_rel_path: str, source_saved_name: str) -> Tuple[str, str, str]:
    rel = (folder_rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    if len(parts) != UPLOAD_ALLOWED_DEPTH:
        raise RuntimeError("Lv5フォルダではないためMarkdown保存先を決定できません。")
    if len(parts) < 5 or parts[3] != "元データ":
        raise RuntimeError("元データ配下のLv5フォルダではないためMarkdown保存先を決定できません。")

    md_parts = list(parts)
    md_parts[3] = "マークダウン形式"
    md_dir_rel = "/".join(md_parts)
    md_dir_abs = resolve_explorer_path(EXPLORER_ROOT, md_dir_rel)

    base_name = os.path.splitext(sanitize_upload_filename(source_saved_name))[0]
    if not base_name:
        raise RuntimeError("Markdownファイル名を決定できません。")

    md_name = f"{base_name}.md"
    md_abs_path = os.path.join(md_dir_abs, md_name)
    md_rel_path = make_rel_from_root(md_abs_path, EXPLORER_ROOT)
    return md_abs_path, md_rel_path, md_name


def is_ondemand_source_folder_rel(rel_path: str) -> bool:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    return len(parts) == UPLOAD_ALLOWED_DEPTH and len(parts) >= 5 and parts[3] == "元データ"


def strip_upload_timestamp_prefix(filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        return ""
    m = re.match(r"^\d{8}_\d{6}_(.+)$", safe)
    if m:
        return sanitize_upload_filename(m.group(1))
    return safe


def build_source_signature(source_abs_path: str, source_rel_path: str) -> str:
    rel = normalize_name_key((source_rel_path or "").replace("\\", "/"))
    try:
        st = os.stat(source_abs_path)
        size = int(st.st_size or 0)
        mtime_ns = int(getattr(st, "st_mtime_ns", int(float(st.st_mtime or 0) * 1_000_000_000)))
        return f"{rel}::{size}:{mtime_ns}"
    except Exception:
        return rel


def delete_ondemand_artifacts(
    folder_rel_path: str,
    source_abs_path: str,
    source_saved_name: str = "",
    source_original_name: str = "",
    require_dify_health: bool = False,
    tolerant: bool = False,
) -> Dict[str, Any]:
    folder_rel = str(folder_rel_path or "").strip().replace("\\", "/").strip("/")
    source_abs = str(source_abs_path or "").strip()
    saved_name = sanitize_upload_filename(source_saved_name) or sanitize_upload_filename(os.path.basename(source_abs)) or os.path.basename(source_abs)
    original_name = sanitize_upload_filename(source_original_name) or strip_upload_timestamp_prefix(saved_name) or sanitize_upload_filename(saved_name)
    if not original_name:
        raise RuntimeError("削除対象ファイル名を判定できません。")

    errors: List[str] = []

    def fail(message: str) -> None:
        msg = safe_err(str(message))
        if tolerant:
            errors.append(msg)
            return
        raise RuntimeError(msg)

    markdown_abs_path = ""
    markdown_rel_path = ""
    markdown_name = ""
    markdown_deleted_files = []
    markdown_deleted = False
    try:
        markdown_abs_path, markdown_rel_path, markdown_name = build_ondemand_markdown_path(folder_rel, saved_name)
        if os.path.isfile(markdown_abs_path):
            try:
                os.remove(markdown_abs_path)
                markdown_deleted = True
                markdown_deleted_files.append(markdown_name)
            except Exception as e:
                fail(f"Markdown削除失敗 ({markdown_name}): {e}")
    except Exception as e:
        fail(f"Markdown削除準備失敗: {e}")

    source_deleted = False
    try:
        if source_abs and os.path.exists(source_abs):
            os.remove(source_abs)
            source_deleted = True
    except Exception as e:
        fail(f"元ファイル削除失敗: {e}")

    return {
        "folder_rel_path": folder_rel,
        "source_abs_path": source_abs,
        "source_name": saved_name,
        "original_name": original_name,
        "source_deleted": bool(source_deleted),
        "markdown_abs_path": markdown_abs_path,
        "markdown_rel_path": markdown_rel_path,
        "markdown_name": markdown_name,
        "markdown_deleted": bool(markdown_deleted),
        "markdown_deleted_files": markdown_deleted_files,
        "errors": errors,
    }


def cleanup_markdown_only(markdown_abs_path: str) -> Dict[str, Any]:
    errors: List[str] = []
    markdown_deleted = False
    if markdown_abs_path and os.path.exists(markdown_abs_path):
        try:
            os.remove(markdown_abs_path)
            markdown_deleted = True
        except Exception as e:
            errors.append(f"Markdown削除失敗: {safe_err(str(e))}")
    return {
        "markdown_deleted": markdown_deleted,
        "errors": errors,
    }


def iter_ondemand_watch_folders(root_dir: str):
    root_abs = resolve_explorer_path(root_dir, "")

    for current_dir, dirnames, _ in os.walk(root_abs):
        depth = path_depth_from_root(current_dir, root_dir)

        if depth >= EXPLORER_MAX_DEPTH:
            dirnames[:] = []
        else:
            next_depth = depth + 1
            dirnames[:] = [d for d in dirnames if matches_explorer_level_rule(next_depth, d)]
            dirnames.sort(key=lambda x: x.lower())

        if depth == UPLOAD_ALLOWED_DEPTH:
            rel = make_rel_from_root(current_dir, root_dir)
            if is_ondemand_source_folder_rel(rel):
                yield current_dir, rel


def list_ondemand_source_files(folder_abs_path: str, root_dir: str) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for name in sorted(os.listdir(folder_abs_path), key=lambda x: x.lower()):
        abs_path = os.path.join(folder_abs_path, name)
        if not os.path.isfile(abs_path):
            continue

        ext = os.path.splitext(name)[1].lower()
        if ext not in ALLOWED_EXTS:
            continue

        rel_path = make_rel_from_root(abs_path, root_dir)
        original_name = strip_upload_timestamp_prefix(name) or sanitize_upload_filename(name)
        if not original_name:
            continue

        out.append({
            "source_abs_path": abs_path,
            "source_rel_path": rel_path,
            "source_saved_name": name,
            "source_original_name": original_name,
            "source_signature": build_source_signature(abs_path, rel_path),
        })
    return out


def build_ondemand_doc_key(dataset_id: str, dataset_name: str, markdown_name: str) -> str:
    ds_key = normalize_name_key(dataset_id or dataset_name)
    md_key = normalize_name_key(markdown_name)
    if not ds_key or not md_key:
        return ""
    return f"{ds_key}::{md_key}"